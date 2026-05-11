#!/usr/bin/env python3
"""Génère le PPTX depuis le HTML de la présentation système de design agentique."""

import re, io, base64
from pathlib import Path
from bs4 import BeautifulSoup, Tag, NavigableString
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import cairosvg
from PIL import Image

DIR  = Path(__file__).parent
HTML = DIR / "présentation-systèmes-design-agentique.html"
OUT  = DIR / "présentation-systèmes-design-agentique.pptx"

def px(n): return Emu(round(n * 9525))
W, H = px(1280), px(720)

def rgb(h):
    h = h.lstrip('#')
    return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

NAVY=rgb('060D1A'); NAVY2=rgb('0A1628'); CYAN=rgb('14D9C8'); BLUE=rgb('2B7FFF')
WHITE=rgb('EEF3FF'); GREY0=rgb('B8CCEE'); GREY=rgb('7A90B8')
AMBER=rgb('ED6B41'); RED=rgb('FF5555'); RED_LT=rgb('FFB0B0'); PURPLE=rgb('9B7FFF')

CSS_VARS = {
    'var(--navy)':'#060D1A','var(--navy2)':'#0A1628','var(--navy3)':'#0F2040',
    'var(--cyan)':'#14D9C8','var(--cyan2)':'#0FBFB0','var(--blue)':'#2B7FFF',
    'var(--white)':'#EEF3FF','var(--grey0)':'#B8CCEE','var(--grey)':'#7A90B8',
    'var(--grey2)':'#4A5E80','var(--amber)':'#ED6B41','var(--red)':'#FF5555',
    'var(--purple)':'#9B7FFF','var(--orange-ramq)':'#ED6B41',
    'var(--navy0)':'#020810','var(--blue-light)':'#7AB8FF',
    'var(--red-light)':'#FFB0B0','var(--purple-light)':'#C5AEFF',
    'var(--blue-mid)':'#78A0DC','var(--cyan-dark)':'#17A79A',
    'var(--gold)':'#F5A623','var(--gold-light)':'#F5C060','var(--blue-pale)':'#C0CDE8',
}
def resolve(s):
    for v, c in CSS_VARS.items(): s = s.replace(v, c)
    return s

F_ATK='Atkinson Hyperlegible'; F_DM='DM Sans'

# ── Shape / text helpers ──────────────────────────────────────────────────────

def set_bg(slide, color=NAVY):
    bg=slide.background; f=bg.fill; f.solid(); f.fore_color.rgb=color

def rect(slide, x, y, w, h, fill=None, line=None):
    s=slide.shapes.add_shape(1, x, y, w, h)
    if fill: s.fill.solid(); s.fill.fore_color.rgb=fill
    else: s.fill.background()
    if line: s.line.color.rgb=line
    else: s.line.fill.background()
    return s

def accent_bar(slide):
    rect(slide, px(0), px(0), px(5), H, fill=CYAN)

def txt(slide, text, x, y, w, h, font=F_DM, size=13, color=WHITE,
        bold=False, italic=False, align=PP_ALIGN.LEFT):
    tb=slide.shapes.add_textbox(x, y, w, h); tb.word_wrap=True
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=align
    run=p.add_run(); run.text=text
    run.font.name=font; run.font.size=Pt(size); run.font.color.rgb=color
    run.font.bold=bold; run.font.italic=italic
    return tb

def label(slide, text, x, y, w=px(700)):
    txt(slide, text.upper(), x, y, w, px(20), font=F_DM, size=9, color=GREY0)

def heading(slide, soup_elem, x, y, w, size=32):
    """Heading with <br>→new para, <em>→cyan, <strong>→white bold."""
    tb=slide.shapes.add_textbox(x, y, w, px(110)); tb.word_wrap=True
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]; p.alignment=PP_ALIGN.LEFT
    for child in soup_elem.children:
        if isinstance(child, NavigableString):
            t=str(child).replace('\xa0',' ')
            if not t: continue
            run=p.add_run(); run.text=t
            run.font.name=F_ATK; run.font.size=Pt(size)
            run.font.color.rgb=WHITE; run.font.bold=True
        elif isinstance(child, Tag):
            if child.name=='br':
                p=tf.add_paragraph(); p.alignment=PP_ALIGN.LEFT
            elif child.name=='em':
                run=p.add_run(); run.text=child.get_text()
                run.font.name=F_ATK; run.font.size=Pt(size)
                sty=child.get('style','')
                run.font.color.rgb=AMBER if 'amber' in sty else CYAN
                run.font.bold=True
            elif child.name=='strong':
                run=p.add_run(); run.text=child.get_text()
                run.font.name=F_ATK; run.font.size=Pt(size)
                run.font.color.rgb=WHITE; run.font.bold=True
    return tb

def rich_para(slide, soup_elem, x, y, w, h, size=13, color=GREY0):
    """Para with <strong>→white bold, <em>→cyan/amber."""
    tb=slide.shapes.add_textbox(x, y, w, h); tb.word_wrap=True
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    for child in soup_elem.children:
        if isinstance(child, NavigableString):
            t=str(child).replace('\xa0',' ')
            if not t: continue
            run=p.add_run(); run.text=t
            run.font.name=F_DM; run.font.size=Pt(size); run.font.color.rgb=color
        elif isinstance(child, Tag):
            if child.name=='strong':
                run=p.add_run(); run.text=child.get_text()
                run.font.name=F_DM; run.font.size=Pt(size)
                sty=child.get('style','')
                c2=CYAN if 'cyan' in sty else WHITE
                run.font.color.rgb=c2; run.font.bold=True
            elif child.name=='em':
                run=p.add_run(); run.text=child.get_text()
                run.font.name=F_DM; run.font.size=Pt(size); run.font.color.rgb=AMBER
            elif child.name=='br':
                p=tf.add_paragraph()
            else:
                run=p.add_run(); run.text=child.get_text()
                run.font.name=F_DM; run.font.size=Pt(size); run.font.color.rgb=color
    return tb

def big_statement(slide, soup_elem, x, y, w, size=20):
    """Big statement text with em→cyan support."""
    src_el=soup_elem.find(class_='src')
    src_t=src_el.get_text(strip=True) if src_el else ''
    if src_el: src_el.extract()
    tb=slide.shapes.add_textbox(x, y, w, px(130)); tb.word_wrap=True
    tf=tb.text_frame; tf.word_wrap=True
    p=tf.paragraphs[0]
    for child in soup_elem.children:
        if isinstance(child, NavigableString):
            t=str(child).replace('\xa0',' ')
            if not t.strip(): continue
            run=p.add_run(); run.text=t
            run.font.name=F_ATK; run.font.size=Pt(size)
            run.font.color.rgb=WHITE; run.font.bold=True
        elif isinstance(child, Tag):
            if child.name=='em':
                run=p.add_run(); run.text=child.get_text()
                run.font.name=F_ATK; run.font.size=Pt(size)
                run.font.color.rgb=CYAN; run.font.bold=True
            elif child.name=='br':
                p=tf.add_paragraph()
            elif child.name=='strong':
                run=p.add_run(); run.text=child.get_text()
                run.font.name=F_ATK; run.font.size=Pt(size)
                run.font.color.rgb=WHITE; run.font.bold=True
    return src_t

def quote_box(slide, text, x, y, w):
    rect(slide, x, y, px(4), px(68), fill=AMBER)
    txt(slide, text[:180], x+px(16), y, w-px(16), px(68),
        size=11, color=WHITE, italic=True)
    return px(78)

def img_from_b64(slide, data_url, x, y, w, h):
    m=re.match(r'data:image/(\w+);base64,(.+)', data_url, re.DOTALL)
    if not m: return
    img_bytes=base64.b64decode(m.group(2))
    img=Image.open(io.BytesIO(img_bytes))
    if img.mode not in ('RGB',): img=img.convert('RGB')
    buf=io.BytesIO(); img.save(buf,'JPEG',quality=82); buf.seek(0)
    slide.shapes.add_picture(buf, x, y, w, h)

def img_from_bytes(slide, png_bytes, x, y, w, h):
    slide.shapes.add_picture(io.BytesIO(png_bytes), x, y, w, h)

def svg_to_png(svg_str, out_w=None, out_h=None):
    s=resolve(svg_str)
    s=re.sub(r'(<svg[^>]*>)',r'\1<rect width="100%" height="100%" fill="#060D1A"/>',s,count=1)
    kw={}
    if out_w: kw['output_width']=int(out_w)
    if out_h: kw['output_height']=int(out_h)
    try: return cairosvg.svg2png(bytestring=s.encode('utf-8'), **kw)
    except Exception as e: print(f"  SVG error: {e}"); return None

# ── Parse ─────────────────────────────────────────────────────────────────────
print("Lecture…")
soup=BeautifulSoup(open(HTML,encoding='utf-8').read(),'lxml')
slides_html=soup.find_all('div',class_='slide')
print(f"  {len(slides_html)} slides")

prs=Presentation(); prs.slide_width=W; prs.slide_height=H
blank=prs.slide_layouts[6]

# ── Slides ────────────────────────────────────────────────────────────────────
for i, sh in enumerate(slides_html):
    sid=sh.get('id',f's{i+1}')
    print(f"  {sid}…", end=' ')
    sl=prs.slides.add_slide(blank)
    set_bg(sl); accent_bar(sl)

    col_full=sh.find('div',class_='col-full')
    lay_spl =sh.find('div',class_='layout-split')

    # ── COL-FULL ──────────────────────────────────────────────────────────────
    if col_full:
        lx,ly,lw=px(72),px(52),W-px(144); cy=ly
        lbl=col_full.find(class_='label')
        hdr=col_full.find(['h1','h2','h3'])
        if lbl: label(sl, lbl.get_text(strip=True), lx, cy); cy+=px(22)
        if hdr:
            sz=28 if 'h2-sm' in hdr.get('class',[]) else 34
            heading(sl, hdr, lx, cy, lw, size=sz); cy+=px(68)

        if sid=='s2':
            cols=col_full.find_all('div',style=lambda s:s and '56px' in s)
            for idx,nd in enumerate(cols[:4]):
                cx=lx+idx*px(290)
                txt(sl, nd.get_text(strip=True), cx, cy, px(270), px(54),
                    font=F_ATK, size=30, color=rgb('14D9C880'), bold=True)
                par=nd.find_parent('div'); paras=par.find_all('p') if par else []
                if paras: rich_para(sl, paras[0], cx, cy+px(56), px(270), px(44), size=12, color=WHITE)
                if len(paras)>1: txt(sl, paras[1].get_text(strip=True), cx, cy+px(104), px(270), px(36), size=10, color=GREY)

        elif sid=='s5':
            cards=col_full.find_all(class_='card-blue')
            pos=[(lx,cy),(lx+px(630),cy),(lx,cy+px(148)),(lx+px(630),cy+px(148))]
            for idx,card in enumerate(cards[:4]):
                cx,ccy=pos[idx]; cw,ch=px(620),px(138)
                rect(sl,cx,ccy,cw,ch,fill=rgb('3421CB40'),line=rgb('FFFFFF14'))
                em=card.find('div',style=lambda s:s and '32px' in s)
                ti=card.find('div',style=lambda s:s and '18px' in s)
                pd=card.find('p')
                if em: txt(sl, em.get_text(strip=True), cx+px(18), ccy+px(18), px(44), px(44), size=22)
                if ti: txt(sl, ti.get_text(strip=True), cx+px(68), ccy+px(18), cw-px(88), px(28), font=F_ATK, size=14, color=WHITE, bold=True)
                if pd: rich_para(sl, pd, cx+px(68), ccy+px(50), cw-px(88), px(72), size=11, color=GREY0)

        elif sid=='s7':
            rows=[
                ('📄','Guide de style','Conventions visuelles ou éditoriales','Artefact partiel','Utile, mais ne couvre pas les règles'),
                ('🎨','Charte graphique','Identité visuelle, logo, couleurs','Intrant, pas le système','Alimente les fondations, mais ne les remplace pas'),
                ('🖥️','Fichier Figma','Outil de travail, même très bien organisé','Le contenant, pas le contenu','Peut héberger des éléments du système'),
                ('🧩','Bibliothèque de composants','Ressource visuelle réutilisable','Une partie du système','Sans règles d\'usage, c\'est une boîte sans mode d\'emploi'),
            ]
            hw=lw//2
            rect(sl,lx,cy,lw,px(34),fill=rgb('FFFFFF14'))
            txt(sl,"CE QU'ON CROIT",lx+px(16),cy+px(9),hw,px(18),font=F_ATK,size=10,color=AMBER,bold=True)
            txt(sl,"CE QUE C'EST VRAIMENT",lx+hw+px(16),cy+px(9),hw,px(18),font=F_ATK,size=10,color=CYAN,bold=True)
            cy+=px(34)
            for row in rows:
                rh=px(102)
                rect(sl,lx,cy,hw,rh,fill=rgb('FF555514'))
                rect(sl,lx+hw,cy,hw,rh,fill=rgb('060D1A'))
                txt(sl,row[0],lx+px(14),cy+px(28),px(28),px(44),size=18)
                txt(sl,row[1],lx+px(50),cy+px(12),hw-px(64),px(28),font=F_ATK,size=13,color=AMBER,bold=True)
                txt(sl,row[2],lx+px(50),cy+px(44),hw-px(64),px(42),size=10,color=GREY)
                txt(sl,row[3],lx+hw+px(16),cy+px(12),hw-px(32),px(28),font=F_ATK,size=13,color=CYAN,bold=True)
                txt(sl,row[4],lx+hw+px(16),cy+px(44),hw-px(32),px(44),size=10,color=GREY)
                cy+=rh
        print("col-full ✓")

    # ── LAYOUT-SPLIT ──────────────────────────────────────────────────────────
    elif lay_spl:
        col_left =lay_spl.find(class_='col-left')
        col_right=lay_spl.find(class_='col-right')

        right_w=px(520)
        if col_right:
            m=re.search(r'width:(\d+)px',col_right.get('style',''))
            if m: right_w=px(int(m.group(1)))
        right_x=W-right_w

        lx,ly=px(68),px(52); lw=right_x-lx-px(32); cy=ly

        # ─ LEFT ──────────────────────────────────────────────────────────────
        if col_left:
            lbl=col_left.find(class_='label')
            hdr=col_left.find(['h1','h2','h3'])
            if lbl: label(sl, lbl.get_text(strip=True), lx, cy, lw); cy+=px(22)
            if hdr:
                is_h1=hdr.name=='h1'
                sz=44 if is_h1 else (28 if 'h2-sm' in hdr.get('class',[]) else 32)
                heading(sl, hdr, lx, cy, lw, size=sz); cy+=px(130) if is_h1 else px(85)
                if is_h1:
                    dp=col_left.find('p')
                    if dp: txt(sl, dp.get_text(strip=True), lx, cy, lw, px(24), size=12, color=GREY); cy+=px(28)

            for par in col_left.find_all('p',limit=3):
                sty=par.get('style','')
                c=GREY if 'var(--grey)' in sty else GREY0
                sz=11 if '13px' in sty else (12 if '14px' in sty else 13)
                rich_para(sl, par, lx, cy, lw, px(56), size=sz, color=c); cy+=px(60)

            q=col_left.find(class_='quote')
            if q: cy+=quote_box(sl, q.get_text(' ',strip=True), lx, cy, lw)

            # Feature items with icon
            for fi in col_left.find_all('div',style=lambda s:s and
                ('rgba(52,33,203' in s or ('rgba(255,85,85' in s and 'align-items:center' in s))
                and 'border-radius:8px' in s):
                spans=fi.find_all('span',limit=2)
                if len(spans)>=2:
                    em,ct=spans[0].get_text(strip=True),spans[1].get_text(strip=True)
                    rect(sl,lx,cy,lw,px(34),fill=rgb('3421CB40'))
                    txt(sl,f'{em}  {ct}',lx+px(12),cy+px(8),lw-px(24),px(20),size=11,color=WHITE)
                    cy+=px(38)

            # Pills
            pills=col_left.find_all(class_=re.compile(r'\bpill\b'))
            if pills:
                pxx=lx
                for pill in pills:
                    pt=pill.get_text(strip=True); pw=px(len(pt)*6+24)
                    is_red='pill-red' in ' '.join(pill.get('class',[]))
                    fc=rgb('FF555526') if is_red else rgb('2B7FFF26')
                    lc=rgb('FF555540') if is_red else rgb('2B7FFF40')
                    rect(sl,pxx,cy,pw,px(22),fill=fc,line=lc)
                    txt(sl,pt,pxx+px(4),cy+px(3),pw-px(8),px(16),size=9,color=GREY0,align=PP_ALIGN.CENTER)
                    pxx+=pw+px(8)
                cy+=px(30)

            # Numbered cards (s12)
            for nc in col_left.find_all(class_='card',limit=3):
                nd=nc.find('div',style=lambda s:s and '42px' in s)
                tit=nc.find('div',style=lambda s:s and '17px' in s)
                pd=nc.find('p'); ch=px(76)
                rect(sl,lx,cy,lw,ch,fill=rgb('FFFFFF26'),line=rgb('FFFFFF40'))
                if nd: txt(sl,nd.get_text(strip=True),lx+px(10),cy+px(8),px(44),px(58),font=F_ATK,size=26,color=rgb('14D9C880'),bold=True)
                if tit: txt(sl,tit.get_text(strip=True),lx+px(60),cy+px(10),lw-px(72),px(28),font=F_ATK,size=13,color=CYAN,bold=True)
                if pd: rich_para(sl,pd,lx+px(60),cy+px(42),lw-px(72),px(26),size=10)
                cy+=ch+px(8)

            # "Ce que ce n'est pas" (s14)
            for ni in col_left.find_all('div',style=lambda s:s and 'rgba(255,85,85,0.15)' in s and 'border-radius:10px' in s):
                spans=ni.find_all('span',limit=2)
                if len(spans)>=2:
                    ct=spans[1].get_text(strip=True)
                    rect(sl,lx,cy,lw,px(34),fill=rgb('FF555526'),line=rgb('FF555533'))
                    txt(sl,f'✗  {ct}',lx+px(12),cy+px(8),lw-px(24),px(20),size=11,color=RED_LT)
                    cy+=px(38)
            pi=col_left.find('div',style=lambda s:s and 'rgba(20,217,200,0.07)' in s)
            if pi:
                spans=pi.find_all('span',limit=2)
                if len(spans)>=2:
                    ct=spans[1].get_text(strip=True)
                    rect(sl,lx,cy,lw,px(34),fill=rgb('14D9C812'),line=rgb('14D9C840'))
                    txt(sl,f'✓  {ct}',lx+px(12),cy+px(8),lw-px(24),px(20),size=11,color=CYAN)

            # Definition box (s11)
            def_box=col_left.find('div',style=lambda s:s and 'rgba(20,217,200,0.07)' in s and 'border-radius:14px' in s)
            if def_box:
                dtxt=def_box.get_text(' ',strip=True)
                rect(sl,lx,cy,lw,px(72),fill=rgb('14D9C812'),line=rgb('14D9C840'))
                rich_para(sl,def_box,lx+px(16),cy+px(10),lw-px(32),px(54),size=12,color=WHITE)
                cy+=px(80)

            # Big statement on left (s11)
            bs_left=col_left.find(class_='big-statement')
            if bs_left:
                heading(sl,bs_left,lx,cy,lw,size=16); cy+=px(60)

            # Cards cyan (s13)
            for cc in col_left.find_all(class_='card',limit=4):
                if 'card-cyan' in ' '.join(cc.get('class',[])):
                    tit=cc.find('strong'); pt=cc.get_text(' ',strip=True)
                    rect(sl,lx,cy,lw,px(60),fill=rgb('14D9C819'),line=rgb('14D9C840'))
                    if tit: txt(sl,tit.get_text(strip=True),lx+px(14),cy+px(8),lw-px(28),px(22),font=F_ATK,size=11,color=CYAN,bold=True)
                    txt(sl,pt,lx+px(14),cy+px(30),lw-px(28),px(24),size=10,color=GREY0)
                    cy+=px(66)

        # ─ RIGHT ─────────────────────────────────────────────────────────────
        if col_right:
            rsty=col_right.get('style','')
            # BG color
            bgc=NAVY2
            if 'rgba(43,127,255,0.05)' in rsty: bgc=rgb('2B7FFF0D')
            elif 'rgba(20,217,200,0.04)' in rsty: bgc=rgb('14D9C80A')
            elif 'rgba(0,0,0' in rsty: bgc=rgb('060D1A')
            rect(sl,right_x,px(0),right_w,H,fill=bgc)

            # Image from <img>
            img_tag=col_right.find('img')
            if img_tag:
                src=img_tag.get('src','')
                if src.startswith('data:image'):
                    img_from_b64(sl, src, right_x, px(0), right_w, H)
                    print("img ✓", end=' ')

            # Background image from CSS
            elif 'url(data:image' in rsty:
                m=re.search(r'url\((data:image/\w+;base64,[^)]+)\)',rsty)
                if m:
                    img_from_b64(sl, m.group(1), right_x, px(0), right_w, H)
                    # Semi-transparent navy overlay
                    ov=rect(sl, right_x, px(0), right_w, H, fill=NAVY)
                    # Set 65% transparency on the overlay shape
                    from lxml import etree as ET
                    from pptx.oxml.ns import qn
                    sp=ov._element
                    for solidFill in sp.iter(qn('a:solidFill')):
                        for srgb in solidFill.iter(qn('a:srgbClr')):
                            alpha=ET.SubElement(srgb, qn('a:alpha'))
                            alpha.set('val','40000')
                    print("bg-img ✓", end=' ')
                # Right column text on top
                rx=right_x+px(48); rw=right_w-px(96); ry=px(235)
                rh2=col_right.find('h2')
                if rh2: heading(sl,rh2,rx,ry,rw,size=24); ry+=px(95)
                fp=col_right.find('p')
                if fp: rich_para(sl,fp,rx,ry,rw,px(80),size=11,color=GREY0)

            # SVG diagram
            elif col_right.find('svg') and not img_tag:
                svg_tag=col_right.find('svg')
                vb=svg_tag.get('viewBox','0 0 500 380'); vp=vb.split()
                svw,svh=(float(vp[2]),float(vp[3])) if len(vp)==4 else (500,380)
                tw=int((right_w-px(64))/9525); th=int(tw*svh/svw)
                png=svg_to_png(str(svg_tag),tw,th)
                if png:
                    dw=right_w-px(64); dh=Emu(int(dw*svh/svw))
                    img_from_bytes(sl,png,right_x+(right_w-dw)//2,(H-dh)//2,dw,dh)
                    print(f"svg→png {tw}×{th} ✓", end=' ')

            # Text content
            elif not img_tag and 'url(data:image' not in rsty:
                rx=right_x+px(48); rw=right_w-px(96); ry=px(140)
                bs=col_right.find(class_='big-statement')
                if bs:
                    src_t=big_statement(sl, bs, rx, ry, rw, size=19)
                    ry+=px(140)
                    if src_t: txt(sl,src_t,rx,ry,rw,px(20),size=9,color=GREY,italic=True); ry+=px(24)
                hr=col_right.find('div',style=lambda s:s and 'height:1px' in s)
                if hr: rect(sl,rx,ry,rw,px(1),fill=rgb('FFFFFF28')); ry+=px(16)
                for rp in col_right.find_all('p',limit=3):
                    t=rp.get_text(' ',strip=True)
                    if not t: continue
                    sty=rp.get('style',''); c=GREY0 if ('grey0' in sty or 'B8CCEE' in sty) else GREY
                    sz=11 if '13px' in sty else (12 if '14px' in sty else 13)
                    rich_para(sl,rp,rx,ry,rw,px(60),size=sz,color=c); ry+=px(66)
                # Token comparison (s10)
                for ti in col_right.find_all('div',style=lambda s:s and
                    ('rgba(255,85,85,0.1)' in s or 'rgba(20,217,200,0.1)' in s)
                    and 'border-radius:10px' in s):
                    t=ti.get_text(strip=True); is_red='rgba(255,85,85' in ti.get('style','')
                    rect(sl,rx,ry,rw,px(40),fill=rgb('FF555519') if is_red else rgb('14D9C819'))
                    txt(sl,t,rx+px(12),ry+px(10),rw-px(24),px(22),font=F_ATK,size=14,
                        color=AMBER if is_red else CYAN,bold=True,align=PP_ALIGN.CENTER)
                    ry+=px(46)
                # Small label above/below tokens
                for li in col_right.find_all('p',style=lambda s:s and '11px' in s and 'uppercase' in s):
                    t=li.get_text(strip=True); c=CYAN if 'var(--cyan)' in li.get('style','') else GREY
                    txt(sl,t,rx,ry,rw,px(18),size=8,color=c,align=PP_ALIGN.CENTER); ry+=px(22)
                # Right cards (s14)
                for rc in col_right.find_all('div',style=lambda s:s and 'rgba(52,33,203,0.25)' in s and 'border-radius:12px' in s):
                    tit=rc.find('div',style=lambda s:s and '15px' in s); pd=rc.find('p')
                    rect(sl,rx,ry,rw,px(70),fill=rgb('3421CB40'))
                    if tit: txt(sl,tit.get_text(strip=True),rx+px(14),ry+px(8),rw-px(28),px(26),font=F_ATK,size=12,color=WHITE,bold=True)
                    if pd: txt(sl,pd.get_text(strip=True),rx+px(14),ry+px(36),rw-px(28),px(28),size=10,color=GREY)
                    ry+=px(78)
                # s15 right heading
                rh2=col_right.find('h2')
                if rh2: heading(sl,rh2,rx,ry,rw,size=24); ry+=px(95)
        print("✓")

print(f"\nSauvegarde → {OUT}")
prs.save(str(OUT))
print("✓ Terminé.")
